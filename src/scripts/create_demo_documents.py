"""
데모용 샘플 문서 생성 스크립트

품질검사 매뉴얼(PPTX)과 프로젝트 가이드(TXT) 파일을 생성합니다.
생성된 파일은 data/demo_documents/ 디렉터리에 저장됩니다.

Usage:
    python -m src.scripts.create_demo_documents
    또는
    python src/scripts/create_demo_documents.py
"""

from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# ---------------------------------------------------------------------------
# 경로 설정
# ---------------------------------------------------------------------------
PROJECT_ROOT = Path(__file__).resolve().parents[2]
OUTPUT_DIR = PROJECT_ROOT / "data" / "demo_documents"


# ===========================================================================
# 유틸리티 헬퍼
# ===========================================================================

def _set_slide_bg(slide, rgb: RGBColor) -> None:
    """슬라이드 배경색을 단색으로 설정합니다."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _add_textbox(
    slide,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    text: str,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = RGBColor(0x33, 0x33, 0x33),
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
    font_name: str = "맑은 고딕",
):
    """텍스트 박스를 추가하고 기본 서식을 적용합니다."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txbox


def _add_bullet_list(
    slide,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    items: list[str],
    font_size: int = 16,
    color: RGBColor = RGBColor(0x33, 0x33, 0x33),
    font_name: str = "맑은 고딕",
    line_spacing: float = 1.5,
):
    """불릿 리스트(bullet list) 텍스트 박스를 추가합니다."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True

    for idx, item in enumerate(items):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1))
        p.level = 0
    return txbox


def _add_table(
    slide,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    rows: list[list[str]],
    header_color: RGBColor = RGBColor(0x2E, 0x56, 0x8C),
    font_name: str = "맑은 고딕",
):
    """표(table)를 추가하고 서식을 적용합니다."""
    num_rows = len(rows)
    num_cols = len(rows[0]) if rows else 0
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.name = font_name
                paragraph.alignment = PP_ALIGN.CENTER

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            if row_idx == 0:
                # 헤더 행 스타일
                cell_fill = cell.fill
                cell_fill.solid()
                cell_fill.fore_color.rgb = header_color
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    paragraph.font.bold = True
            else:
                cell_fill = cell.fill
                cell_fill.solid()
                if row_idx % 2 == 1:
                    cell_fill.fore_color.rgb = RGBColor(0xE8, 0xEE, 0xF5)
                else:
                    cell_fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    return table_shape


# ===========================================================================
# 1. PPTX 생성 — 품질검사 매뉴얼
# ===========================================================================

def create_quality_inspection_pptx(output_dir: Path | str | None = None) -> Path:
    """품질검사 매뉴얼 PPTX 파일을 생성합니다.

    Args:
        output_dir: 파일이 저장될 디렉터리. 기본값은 ``data/demo_documents/``.

    Returns:
        생성된 파일의 ``Path`` 객체.
    """
    output_dir = Path(output_dir) if output_dir else OUTPUT_DIR
    output_dir.mkdir(parents=True, exist_ok=True)
    filepath = output_dir / "품질검사_매뉴얼.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    ACCENT_BLUE = RGBColor(0x2E, 0x56, 0x8C)
    DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
    LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    # ----- 슬라이드 1: 표지 -----
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _set_slide_bg(slide1, ACCENT_BLUE)

    _add_textbox(
        slide1,
        left=Inches(1.5),
        top=Inches(2.0),
        width=Inches(10),
        height=Inches(1.5),
        text="품질검사 매뉴얼 2026",
        font_size=44,
        bold=True,
        color=WHITE,
        alignment=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide1,
        left=Inches(1.5),
        top=Inches(3.8),
        width=Inches(10),
        height=Inches(1.0),
        text="품질관리팀",
        font_size=28,
        bold=False,
        color=RGBColor(0xBB, 0xD5, 0xED),
        alignment=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide1,
        left=Inches(1.5),
        top=Inches(5.2),
        width=Inches(10),
        height=Inches(0.6),
        text="문서 버전 1.0  |  발행일: 2026-01-15  |  대외비",
        font_size=14,
        color=RGBColor(0x8A, 0xB4, 0xD8),
        alignment=PP_ALIGN.CENTER,
    )

    # ----- 슬라이드 2: 검사원 공통 절차 -----
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide2, LIGHT_BG)

    _add_textbox(
        slide2,
        left=Inches(0.8),
        top=Inches(0.4),
        width=Inches(11),
        height=Inches(0.8),
        text="검사원 공통 절차",
        font_size=32,
        bold=True,
        color=ACCENT_BLUE,
    )

    procedures = [
        "1. 검사 준비  —  작업 지시서 확인 및 검사 대상 로트(Lot) 정보 파악",
        "2. 장비 점검  —  측정 장비 교정 상태 확인(교정 유효기간 준수 여부)",
        "3. 샘플링  —  KS Q ISO 2859-1 기준에 따른 샘플 추출",
        "4. 검사 수행  —  검사 기준서에 명시된 항목별 검사 실시",
        "5. 기록 작성  —  검사 결과를 품질관리 시스템(QMS)에 즉시 입력",
        "6. 판정 및 보고  —  합격/불합격 판정 후 관리자 승인 요청",
        "7. 부적합 처리  —  불합격 시 부적합 보고서(NCR) 발행 및 격리 조치",
    ]

    _add_bullet_list(
        slide2,
        left=Inches(1.0),
        top=Inches(1.5),
        width=Inches(11),
        height=Inches(5.5),
        items=procedures,
        font_size=17,
        color=DARK_TEXT,
    )

    # ----- 슬라이드 3: 외관 검사 기준 -----
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide3, LIGHT_BG)

    _add_textbox(
        slide3,
        left=Inches(0.8),
        top=Inches(0.4),
        width=Inches(11),
        height=Inches(0.8),
        text="외관 검사 기준",
        font_size=32,
        bold=True,
        color=ACCENT_BLUE,
    )

    _add_textbox(
        slide3,
        left=Inches(0.8),
        top=Inches(1.3),
        width=Inches(11),
        height=Inches(0.6),
        text="아래 항목에 대해 육안 및 확대경(10x) 검사를 실시합니다.",
        font_size=16,
        color=DARK_TEXT,
    )

    inspection_table = [
        ["검사 항목", "판정 기준", "검사 방법", "허용 한도"],
        ["스크래치", "표면 길이 0.5mm 이하", "육안 / 확대경(10x)", "A면 불허 / B면 1개 이하"],
        ["변색", "색차 ΔE ≤ 1.5", "색차계(Spectrophotometer)", "전 면 동일 기준 적용"],
        ["치수", "도면 공차 ±0.05mm", "버니어 캘리퍼스 / CMM", "공차 초과 시 불합격"],
        ["이물질", "부착 이물 없을 것", "육안 / UV 조명", "전수 검사 기준"],
        ["용접 상태", "기공·크랙 없을 것", "육안 / X-ray(필요시)", "KS B 0845 기준 적용"],
    ]

    _add_table(
        slide3,
        left=Inches(0.8),
        top=Inches(2.1),
        width=Inches(11.5),
        height=Inches(4.5),
        rows=inspection_table,
        header_color=ACCENT_BLUE,
    )

    # ----- 슬라이드 4: 불량 유형 분류 -----
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide4, LIGHT_BG)

    _add_textbox(
        slide4,
        left=Inches(0.8),
        top=Inches(0.4),
        width=Inches(11),
        height=Inches(0.8),
        text="불량 유형 분류",
        font_size=32,
        bold=True,
        color=ACCENT_BLUE,
    )

    _add_textbox(
        slide4,
        left=Inches(0.8),
        top=Inches(1.3),
        width=Inches(11),
        height=Inches(0.6),
        text="결함 등급에 따라 조치 수준이 달라집니다. 아래 분류표를 참고하십시오.",
        font_size=16,
        color=DARK_TEXT,
    )

    defect_table = [
        ["등급", "분류명", "정의", "예시", "조치 사항"],
        [
            "A등급",
            "치명결함\n(Critical)",
            "사용자 안전에 직접적 위험을\n초래하거나 법적 규제를 위반",
            "절연 파괴, 날카로운 모서리,\n유해물질 기준 초과",
            "즉시 생산 중단\n전수 검사 실시\n원인 분석(8D) 필수",
        ],
        [
            "B등급",
            "중결함\n(Major)",
            "제품 기능 또는 성능 저하를\n유발하여 고객 불만 가능",
            "치수 공차 초과, 동작 불량,\n외관 A면 스크래치",
            "해당 로트 격리\n샘플링 강화 검사\n시정 조치 보고서 발행",
        ],
        [
            "C등급",
            "경결함\n(Minor)",
            "제품 사용에 큰 영향 없으나\n미관 또는 마감 품질 저하",
            "미세 변색, B면 경미한 스크래치,\n포장 라벨 위치 미세 어긋남",
            "기록 후 출하 허용\n(허용 한도 내)\n주간 품질 회의 보고",
        ],
    ]

    _add_table(
        slide4,
        left=Inches(0.5),
        top=Inches(2.1),
        width=Inches(12.3),
        height=Inches(4.8),
        rows=defect_table,
        header_color=ACCENT_BLUE,
    )

    # ----- 슬라이드 5: 검사 보고서 작성법 -----
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide5, LIGHT_BG)

    _add_textbox(
        slide5,
        left=Inches(0.8),
        top=Inches(0.4),
        width=Inches(11),
        height=Inches(0.8),
        text="검사 보고서 작성법",
        font_size=32,
        bold=True,
        color=ACCENT_BLUE,
    )

    report_items = [
        "■ 일일보고서 (Daily Report)",
        "   · 작성 시점: 매일 검사 종료 후 30분 이내",
        "   · 포함 항목: 검사 로트 수, 합격/불합격 수량, 주요 불량 유형, 특이사항",
        "   · 제출 대상: 품질관리 파트장",
        "",
        "■ 주간보고서 (Weekly Report)",
        "   · 작성 시점: 매주 금요일 16:00까지",
        "   · 포함 항목: 주간 불량률 추이, Top 3 불량 유형, 시정조치 진행 현황",
        "   · 제출 대상: 품질관리팀장",
        "",
        "■ 월간보고서 (Monthly Report)",
        "   · 작성 시점: 익월 3영업일 이내",
        "   · 포함 항목: 월간 품질 KPI(불량률, 고객 클레임 건수, CPK 지수),",
        "     개선 과제 목록, 공정 능력 분석 결과, 차월 품질 목표",
        "   · 제출 대상: 품질관리 임원 / 경영 회의 보고용",
    ]

    _add_bullet_list(
        slide5,
        left=Inches(1.0),
        top=Inches(1.4),
        width=Inches(11),
        height=Inches(5.5),
        items=report_items,
        font_size=16,
        color=DARK_TEXT,
        line_spacing=1.3,
    )

    # 저장
    prs.save(str(filepath))
    print(f"[OK] PPTX 생성 완료: {filepath}")
    return filepath


# ===========================================================================
# 2. TXT 생성 — 프로젝트 가이드
# ===========================================================================

def create_project_guide_txt(output_dir: Path | str | None = None) -> Path:
    """MemGate 프로젝트 가이드 TXT 파일을 생성합니다.

    Args:
        output_dir: 파일이 저장될 디렉터리. 기본값은 ``data/demo_documents/``.

    Returns:
        생성된 파일의 ``Path`` 객체.
    """
    output_dir = Path(output_dir) if output_dir else OUTPUT_DIR
    output_dir.mkdir(parents=True, exist_ok=True)
    filepath = output_dir / "프로젝트_가이드.txt"

    content = """\
================================================================================
                    MemGate 프로젝트 가이드라인
================================================================================

1. 프로젝트 개요
--------------------------------------------------------------------------------
MemGate는 AI 기반 메모리 에이전트 시스템으로, 사용자와의 대화 이력 및 문서를
지능적으로 관리하고 검색할 수 있는 플랫폼입니다.

핵심 목표:
  - 대화 맥락을 장기 기억으로 저장·관리
  - 다양한 형식의 문서(PDF, PPTX, DOCX, TXT 등) 인덱싱 및 시맨틱 검색
  - 사용자 질문에 대해 관련 기억과 문서를 기반으로 정확한 답변 생성

2. 아키텍처
--------------------------------------------------------------------------------
  ┌─────────────┐     ┌──────────────┐     ┌─────────────────┐
  │  Frontend    │────▶│  FastAPI      │────▶│  Vector Store   │
  │  (React)     │     │  Backend      │     │  (ChromaDB)     │
  └─────────────┘     └──────┬───────┘     └─────────────────┘
                             │
                      ┌──────▼───────┐
                      │  LLM Engine  │
                      │  (OpenAI)    │
                      └──────────────┘

3. 개발 환경 설정
--------------------------------------------------------------------------------
  3.1 필수 요구 사항
    - Python 3.11 이상
    - Node.js 18 이상 (프론트엔드 개발 시)
    - Docker & Docker Compose (배포 환경)

  3.2 로컬 개발 환경 구축
    $ git clone https://github.com/memgate/ai-memory-agent.git
    $ cd ai-memory-agent
    $ uv venv .venv && source .venv/bin/activate
    $ uv pip install -e ".[dev]"
    $ cp .env.example .env   # API 키 등 환경변수 설정
    $ python -m uvicorn src.main:app --reload

4. 코딩 컨벤션
--------------------------------------------------------------------------------
  4.1 Python
    - PEP 8 준수 (black + ruff 자동 포맷팅)
    - 타입 힌트 필수 적용 (from __future__ import annotations)
    - Docstring은 Google 스타일 사용
    - 모든 퍼블릭 함수에 단위 테스트 작성

  4.2 커밋 메시지
    - Conventional Commits 형식 준수
      예: feat: 문서 인덱싱 파이프라인 추가
          fix: 벡터 검색 시 한글 토크나이저 오류 수정
          docs: API 명세서 업데이트

  4.3 브랜치 전략
    - main: 프로덕션 배포 브랜치
    - develop: 통합 개발 브랜치
    - feature/*: 기능 개발 브랜치
    - hotfix/*: 긴급 수정 브랜치

5. 문서 처리 파이프라인
--------------------------------------------------------------------------------
  5.1 지원 파일 형식
    - PDF  → PyMuPDF(fitz) 파서
    - PPTX → python-pptx 파서
    - DOCX → python-docx 파서
    - TXT  → 일반 텍스트 파서
    - CSV  → pandas 기반 파서

  5.2 처리 흐름
    문서 업로드 → 텍스트 추출 → 청크 분할(chunk) → 임베딩 생성 → 벡터 DB 저장

  5.3 청크 분할 규칙
    - 기본 청크 크기: 512 토큰
    - 오버랩: 64 토큰
    - 구분자 우선순위: 문단 > 문장 > 공백

6. API 엔드포인트 요약
--------------------------------------------------------------------------------
  POST   /api/v1/documents/upload    문서 업로드 및 인덱싱
  GET    /api/v1/documents           문서 목록 조회
  GET    /api/v1/documents/{id}      문서 상세 조회
  DELETE /api/v1/documents/{id}      문서 삭제
  POST   /api/v1/search              시맨틱 검색
  POST   /api/v1/chat                대화 (메모리 기반 응답)
  GET    /api/v1/memories            메모리 목록 조회

7. 테스트
--------------------------------------------------------------------------------
  $ pytest tests/ -v --cov=src --cov-report=html
  $ pytest tests/unit/           # 단위 테스트만 실행
  $ pytest tests/integration/    # 통합 테스트만 실행

  - 테스트 커버리지 목표: 80% 이상
  - PR 머지 전 반드시 CI 통과 확인

8. 배포
--------------------------------------------------------------------------------
  8.1 Docker 배포
    $ docker compose -f docker-compose.prod.yml up -d

  8.2 환경 변수 (필수)
    OPENAI_API_KEY      OpenAI API 키
    CHROMA_HOST         ChromaDB 호스트 주소
    CHROMA_PORT         ChromaDB 포트 (기본 8000)
    DATABASE_URL        PostgreSQL 연결 문자열
    SECRET_KEY          JWT 서명 키

9. 팀 연락처
--------------------------------------------------------------------------------
  - 프로젝트 리드: 김영수 (youngsoo.kim@memgate.io)
  - 백엔드 리드: 이정민 (jungmin.lee@memgate.io)
  - AI/ML 리드: 박서연 (seoyeon.park@memgate.io)
  - DevOps: 최동현 (donghyun.choi@memgate.io)

================================================================================
                         문서 끝 — MemGate Team © 2026
================================================================================
"""

    filepath.write_text(content, encoding="utf-8")
    print(f"[OK] TXT 생성 완료: {filepath}")
    return filepath


# ===========================================================================
# 통합 생성 함수
# ===========================================================================

def create_all_demo_documents(output_dir: Path | str | None = None) -> list[Path]:
    """모든 데모 문서를 생성합니다.

    Args:
        output_dir: 파일이 저장될 디렉터리. ``None``이면 기본 경로 사용.

    Returns:
        생성된 파일 경로 목록.
    """
    files: list[Path] = []
    files.append(create_quality_inspection_pptx(output_dir))
    files.append(create_project_guide_txt(output_dir))
    print(f"\n총 {len(files)}개 데모 문서가 생성되었습니다.")
    return files


# ===========================================================================
# CLI 진입점
# ===========================================================================

if __name__ == "__main__":
    create_all_demo_documents()

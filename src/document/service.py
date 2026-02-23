"""문서 서비스"""

import os
import shutil
import subprocess
import uuid
from io import BytesIO
from pathlib import Path
from typing import Any

import aiosqlite
import chardet
from PyPDF2 import PdfReader

from src.document.repository import DocumentRepository
from src.shared.vector_store import upsert_vector, delete_vector, search_vectors

# 슬라이드 이미지 저장 기본 경로
SLIDES_DATA_DIR = Path("data/slides")


class DocumentService:
    """문서 비즈니스 로직"""

    def __init__(self, db: aiosqlite.Connection):
        self.repo = DocumentRepository(db)

    async def upload_document(
        self,
        file_content: bytes,
        filename: str,
        owner_id: str,
        chat_room_id: str | None = None,
    ) -> dict[str, Any]:
        """문서 업로드 → 텍스트 추출 → 청킹 → 임베딩 → 저장"""
        file_type = self._get_file_type(filename)

        doc = await self.repo.create_document(
            name=filename,
            file_type=file_type,
            file_size=len(file_content),
            owner_id=owner_id,
            chat_room_id=chat_room_id,
        )

        try:
            # PPTX 전용 처리
            if file_type == "pptx":
                await self._process_pptx(doc, file_content, owner_id, chat_room_id)
            else:
                await self._process_text_document(doc, file_content, file_type, owner_id, chat_room_id)

        except ValueError:
            raise
        except Exception as e:
            await self.repo.update_document_status(doc["id"], "failed")
            raise e

        return await self.repo.get_document(doc["id"])

    async def _process_text_document(
        self,
        doc: dict[str, Any],
        file_content: bytes,
        file_type: str,
        owner_id: str,
        chat_room_id: str | None,
    ) -> None:
        """PDF/TXT 문서 처리 (기존 로직)"""
        text = self._extract_text(file_content, file_type)
        if not text.strip():
            await self.repo.update_document_status(doc["id"], "failed")
            raise ValueError("문서에서 텍스트를 추출할 수 없습니다")

        chunks = self._chunk_text(text, chunk_size=800, overlap=100)

        embedding_provider = None
        try:
            from src.shared.providers import get_embedding_provider
            embedding_provider = get_embedding_provider()
        except Exception:
            pass

        for i, chunk_text in enumerate(chunks):
            vector_id = str(uuid.uuid4())
            vector = None

            if embedding_provider:
                try:
                    vector = await embedding_provider.embed(chunk_text)
                except Exception as e:
                    print(f"  임베딩 실패 (청크 {i}): {e}")

            chunk_record = await self.repo.create_chunk(doc["id"], chunk_text, i, vector_id)

            try:
                await self.repo.db.execute(
                    "INSERT INTO document_chunks_fts (content, chunk_id, document_id) VALUES (?, ?, ?)",
                    (chunk_text, chunk_record["id"], doc["id"]),
                )
            except Exception:
                pass

            if vector:
                await upsert_vector(vector_id, vector, {
                    "document_id": doc["id"],
                    "chunk_index": i,
                    "chat_room_id": chat_room_id,
                    "owner_id": owner_id,
                    "scope": "document",
                })

        await self.repo.update_document_status(doc["id"], "completed", len(chunks))

        if chat_room_id:
            await self.repo.link_document_to_room(doc["id"], chat_room_id)

    async def _process_pptx(
        self,
        doc: dict[str, Any],
        file_content: bytes,
        owner_id: str,
        chat_room_id: str | None,
    ) -> None:
        """PPTX 문서 처리: 1 slide = 1 chunk"""
        slides_text = self._extract_pptx_slides(file_content)
        if not slides_text:
            await self.repo.update_document_status(doc["id"], "failed")
            raise ValueError("PPTX에서 슬라이드를 추출할 수 없습니다")

        # 슬라이드 이미지 변환
        slide_images = self._convert_slides_to_images(file_content, doc["id"])

        embedding_provider = None
        try:
            from src.shared.providers import get_embedding_provider
            embedding_provider = get_embedding_provider()
        except Exception:
            pass

        for i, slide_text in enumerate(slides_text):
            slide_number = i + 1
            vector_id = str(uuid.uuid4())
            vector = None

            # 슬라이드 이미지 경로
            image_path = slide_images.get(slide_number)

            if embedding_provider and slide_text.strip():
                try:
                    vector = await embedding_provider.embed(slide_text)
                except Exception as e:
                    print(f"  임베딩 실패 (슬라이드 {slide_number}): {e}")

            chunk_record = await self.repo.create_chunk(
                document_id=doc["id"],
                content=slide_text,
                chunk_index=i,
                vector_id=vector_id,
                slide_number=slide_number,
                slide_image_path=image_path,
            )

            # FTS5 인덱스
            try:
                await self.repo.db.execute(
                    "INSERT INTO document_chunks_fts (content, chunk_id, document_id) VALUES (?, ?, ?)",
                    (slide_text, chunk_record["id"], doc["id"]),
                )
            except Exception:
                pass

            if vector:
                await upsert_vector(vector_id, vector, {
                    "document_id": doc["id"],
                    "chunk_index": i,
                    "slide_number": slide_number,
                    "chat_room_id": chat_room_id,
                    "owner_id": owner_id,
                    "scope": "document",
                })

        await self.repo.update_document_status(doc["id"], "completed", len(slides_text))

        if chat_room_id:
            await self.repo.link_document_to_room(doc["id"], chat_room_id)

    def _extract_pptx_slides(self, content: bytes) -> list[str]:
        """python-pptx로 슬라이드별 텍스트 추출 (제목, 본문, 표 포함)"""
        from pptx import Presentation
        from pptx.table import Table

        prs = Presentation(BytesIO(content))
        slides_text = []

        for slide in prs.slides:
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
                if shape.has_table:
                    table: Table = shape.table
                    for row in table.rows:
                        row_texts = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                row_texts.append(cell_text)
                        if row_texts:
                            texts.append(" | ".join(row_texts))

            slides_text.append("\n".join(texts) if texts else "")

        return slides_text

    def _convert_slides_to_images(
        self, file_content: bytes, doc_id: str
    ) -> dict[int, str]:
        """LibreOffice headless → PDF → pdf2image로 슬라이드별 PNG 저장.
        LibreOffice 미설치 시 graceful fallback (빈 dict 반환)."""
        slide_images: dict[int, str] = {}

        # 출력 디렉토리 생성
        output_dir = SLIDES_DATA_DIR / doc_id
        output_dir.mkdir(parents=True, exist_ok=True)

        try:
            import tempfile
            with tempfile.TemporaryDirectory() as tmpdir:
                # PPTX 파일을 임시 디렉토리에 저장
                pptx_path = Path(tmpdir) / "presentation.pptx"
                pptx_path.write_bytes(file_content)

                # LibreOffice headless로 PDF 변환
                pdf_path = Path(tmpdir) / "presentation.pdf"
                result = subprocess.run(
                    [
                        "libreoffice", "--headless", "--convert-to", "pdf",
                        "--outdir", tmpdir, str(pptx_path),
                    ],
                    capture_output=True,
                    timeout=120,
                )

                if result.returncode != 0 or not pdf_path.exists():
                    print(f"LibreOffice PDF 변환 실패: {result.stderr.decode()}")
                    return slide_images

                # pdf2image로 슬라이드별 PNG 생성
                from pdf2image import convert_from_path
                images = convert_from_path(str(pdf_path), dpi=150)

                for i, image in enumerate(images):
                    slide_number = i + 1
                    image_filename = f"slide_{slide_number}.png"
                    image_path = output_dir / image_filename
                    image.save(str(image_path), "PNG")
                    slide_images[slide_number] = str(image_path)

                print(f"슬라이드 이미지 생성 완료: {len(slide_images)}개 ({doc_id})")

        except FileNotFoundError:
            print("LibreOffice가 설치되어 있지 않습니다. 텍스트만 저장합니다.")
        except Exception as e:
            print(f"슬라이드 이미지 변환 실패: {e}. 텍스트만 저장합니다.")

        return slide_images

    def get_slide_image_path(self, doc_id: str, slide_number: int) -> str | None:
        """슬라이드 이미지 파일 경로 반환"""
        image_path = SLIDES_DATA_DIR / doc_id / f"slide_{slide_number}.png"
        if image_path.exists():
            return str(image_path)
        return None

    async def delete_document(self, doc_id: str, user_id: str) -> bool:
        """문서 삭제 (DB + Qdrant + 슬라이드 이미지)"""
        doc = await self.repo.get_document(doc_id)
        if not doc:
            raise ValueError("문서를 찾을 수 없습니다")
        if doc["owner_id"] != user_id:
            raise PermissionError("문서 삭제 권한이 없습니다")

        vector_ids = await self.repo.get_chunk_vector_ids(doc_id)
        for vid in vector_ids:
            await delete_vector(vid)

        # FTS5 인덱스에서도 삭제
        try:
            await self.repo.db.execute(
                "DELETE FROM document_chunks_fts WHERE document_id = ?",
                (doc_id,),
            )
        except Exception:
            pass

        # 슬라이드 이미지 디렉토리 삭제
        slide_dir = SLIDES_DATA_DIR / doc_id
        if slide_dir.exists():
            shutil.rmtree(slide_dir, ignore_errors=True)

        return await self.repo.delete_document(doc_id)

    async def list_documents(
        self,
        owner_id: str | None = None,
        chat_room_id: str | None = None,
    ) -> list[dict[str, Any]]:
        return await self.repo.list_documents(
            owner_id=owner_id, chat_room_id=chat_room_id
        )

    async def get_document_detail(self, doc_id: str) -> dict[str, Any]:
        doc = await self.repo.get_document(doc_id)
        if not doc:
            raise ValueError("문서를 찾을 수 없습니다")
        chunks = await self.repo.get_chunks(doc_id)

        # 슬라이드 이미지 URL 추가
        for chunk in chunks:
            if chunk.get("slide_number"):
                chunk["slide_image_url"] = f"/api/v1/documents/{doc_id}/slides/{chunk['slide_number']}"

        doc["chunks"] = chunks
        return doc

    async def link_to_room(
        self, doc_id: str, room_id: str, user_id: str
    ) -> dict[str, Any]:
        doc = await self.repo.get_document(doc_id)
        if not doc:
            raise ValueError("문서를 찾을 수 없습니다")
        return await self.repo.link_document_to_room(doc_id, room_id)

    async def unlink_from_room(
        self, doc_id: str, room_id: str
    ) -> bool:
        return await self.repo.unlink_document_from_room(doc_id, room_id)

    async def search_documents(
        self, query: str, user_id: str, limit: int = 10
    ) -> list[dict[str, Any]]:
        """사용자의 모든 문서에서 검색"""
        documents = await self.repo.list_documents(owner_id=user_id)
        if not documents:
            return []

        doc_ids = [doc["id"] for doc in documents]

        try:
            from src.shared.providers import get_embedding_provider
            embedding_provider = get_embedding_provider()
            query_vector = await embedding_provider.embed(query)
        except Exception as e:
            print(f"임베딩 실패: {e}")
            return []

        results = await search_vectors(
            query_vector=query_vector,
            limit=limit,
            filter_conditions={
                "scope": "document",
                "document_id": doc_ids,
            },
        )

        enriched = []
        for r in results:
            doc_id = r["payload"].get("document_id")
            chunk_idx = r["payload"].get("chunk_index")
            doc = await self.repo.get_document(doc_id)

            chunks = await self.repo.get_chunks(doc_id)
            chunk_content = ""
            slide_number = None
            slide_image_url = None
            for c in chunks:
                if c["chunk_index"] == chunk_idx:
                    chunk_content = c["content"]
                    slide_number = c.get("slide_number")
                    if slide_number:
                        slide_image_url = f"/api/v1/documents/{doc_id}/slides/{slide_number}"
                    break

            enriched.append({
                "content": chunk_content,
                "score": r["score"],
                "document_name": doc["name"] if doc else "Unknown",
                "document_id": doc_id,
                "chunk_index": chunk_idx,
                "file_type": doc["file_type"] if doc else "unknown",
                "slide_number": slide_number,
                "slide_image_url": slide_image_url,
            })

        return enriched

    async def search_document_chunks(
        self,
        query: str,
        chat_room_id: str,
        limit: int = 5,
    ) -> list[dict[str, Any]]:
        """대화방 연결 문서에서 관련 청크 검색"""
        doc_ids = await self.repo.get_linked_document_ids(chat_room_id)
        if not doc_ids:
            return []

        try:
            from src.shared.providers import get_embedding_provider
            embedding_provider = get_embedding_provider()
            query_vector = await embedding_provider.embed(query)
        except Exception:
            return []

        results = await search_vectors(
            query_vector=query_vector,
            limit=limit,
            filter_conditions={
                "scope": "document",
                "document_id": doc_ids,
            },
        )

        enriched = []
        for r in results:
            doc_id = r["payload"].get("document_id")
            chunk_idx = r["payload"].get("chunk_index")
            doc = await self.repo.get_document(doc_id)

            chunks = await self.repo.get_chunks(doc_id)
            chunk_content = ""
            slide_number = None
            slide_image_url = None
            for c in chunks:
                if c["chunk_index"] == chunk_idx:
                    chunk_content = c["content"]
                    slide_number = c.get("slide_number")
                    if slide_number:
                        slide_image_url = f"/api/v1/documents/{doc_id}/slides/{slide_number}"
                    break

            enriched.append({
                "content": chunk_content,
                "score": r["score"],
                "document_name": doc["name"] if doc else "Unknown",
                "chunk_index": chunk_idx,
                "document_id": doc_id,
                "slide_number": slide_number,
                "slide_image_url": slide_image_url,
            })

        return enriched

    def _get_file_type(self, filename: str) -> str:
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        if ext == "pdf":
            return "pdf"
        elif ext == "txt":
            return "txt"
        elif ext == "pptx":
            return "pptx"
        raise ValueError(f"지원하지 않는 파일 형식입니다. PDF, TXT, PPTX 파일만 지원합니다.")

    def _extract_text(self, content: bytes, file_type: str) -> str:
        if file_type == "pdf":
            reader = PdfReader(BytesIO(content))
            text = ""
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
            return text
        else:
            detected = chardet.detect(content)
            encoding = detected.get("encoding", "utf-8") or "utf-8"
            return content.decode(encoding)

    def _chunk_text(
        self, text: str, chunk_size: int = 800, overlap: int = 100
    ) -> list[str]:
        """문장/문단 경계를 존중하는 시맨틱 청킹"""
        import re

        # 1. 문단 단위로 먼저 분리
        paragraphs = re.split(r'\n\s*\n', text)

        # 2. 각 문단을 문장 단위로 분리
        sentences = []
        for para in paragraphs:
            sents = re.split(r'(?<=[.!?。])\s+', para.strip())
            sentences.extend([s for s in sents if s.strip()])
            sentences.append("")  # 문단 구분자

        # 3. 문장들을 chunk_size 이내로 묶기
        chunks = []
        current_chunk = []
        current_len = 0

        for sent in sentences:
            sent_len = len(sent)
            if current_len + sent_len > chunk_size and current_chunk:
                chunks.append(" ".join(current_chunk).strip())
                # overlap: 마지막 몇 문장을 다음 청크에 포함
                overlap_chunk = []
                overlap_len = 0
                for s in reversed(current_chunk):
                    if overlap_len + len(s) > overlap:
                        break
                    overlap_chunk.insert(0, s)
                    overlap_len += len(s)
                current_chunk = overlap_chunk
                current_len = overlap_len
            current_chunk.append(sent)
            current_len += sent_len

        if current_chunk:
            final = " ".join(current_chunk).strip()
            if final:
                chunks.append(final)

        return [c for c in chunks if c.strip()]
